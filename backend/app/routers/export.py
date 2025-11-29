"""
PowerPoint export router
"""
from flask import Blueprint, jsonify, request, send_file, abort
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import tempfile
import os

bp = Blueprint('export', __name__)


def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    if subtitle and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    return slide


def create_content_slide(prs, title, content):
    """Create a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content
    return slide


def create_chart_slide(prs, title, chart_data, chart_type="bar"):
    """Create a slide with a chart"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    
    # Create chart data
    cd = CategoryChartData()
    cd.categories = chart_data.get("categories", [])
    for series in chart_data.get("series", []):
        cd.add_series(series["name"], series["values"])
    
    # Determine chart type
    xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    if chart_type == "line":
        xl_chart_type = XL_CHART_TYPE.LINE
    elif chart_type == "bar":
        xl_chart_type = XL_CHART_TYPE.BAR_CLUSTERED
    elif chart_type == "pie":
        xl_chart_type = XL_CHART_TYPE.PIE
    
    # Add chart
    x, y, cx, cy = Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)
    slide.shapes.add_chart(xl_chart_type, x, y, cx, cy, cd)
    
    return slide


def generate_business_case_pptx(data):
    """Generate a complete business case presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    project_name = data.get('project_name', 'Business Case')
    description = data.get('description', 'Business Case Analysis')
    financial_data = data.get('financial_data', [])
    assumptions = data.get('assumptions', {})
    
    # Slide 1: Title
    create_title_slide(prs, project_name, description)
    
    # Slide 2: Executive Summary
    create_content_slide(prs, "Executive Summary", 
                        "• Project overview and objectives\n• Key financial highlights\n• Strategic alignment")
    
    # Slide 3: Revenue Projection
    if financial_data:
        years = [str(fd.get('year', '')) for fd in financial_data]
        revenues = [fd.get('revenue', 0) for fd in financial_data]
        create_chart_slide(prs, "Revenue Projection", {
            "categories": years,
            "series": [{"name": "Revenue", "values": revenues}]
        }, "column")
    
    # Slide 4: Cost Analysis
    if financial_data:
        costs = [fd.get('costs', 0) for fd in financial_data]
        op_expenses = [fd.get('operating_expenses', 0) for fd in financial_data]
        create_chart_slide(prs, "Cost Analysis", {
            "categories": years,
            "series": [
                {"name": "Direct Costs", "values": costs},
                {"name": "Operating Expenses", "values": op_expenses}
            ]
        }, "column")
    
    # Slide 5: Profitability
    if financial_data:
        gross_profits = [fd.get('gross_profit', 0) for fd in financial_data]
        ebitda = [fd.get('ebitda', 0) for fd in financial_data]
        create_chart_slide(prs, "Profitability Analysis", {
            "categories": years,
            "series": [
                {"name": "Gross Profit", "values": gross_profits},
                {"name": "EBITDA", "values": ebitda}
            ]
        }, "line")
    
    # Slide 6: Net Income
    if financial_data:
        net_income = [fd.get('net_income', 0) for fd in financial_data]
        create_chart_slide(prs, "Net Income Projection", {
            "categories": years,
            "series": [{"name": "Net Income", "values": net_income}]
        }, "column")
    
    # Slide 7: Key Assumptions
    assumptions_text = "\n".join([f"• {k}: {v}" for k, v in assumptions.items()])
    if not assumptions_text:
        assumptions_text = "• Market growth rate assumptions\n• Cost escalation factors\n• Pricing strategy"
    create_content_slide(prs, "Key Assumptions", assumptions_text)
    
    # Slide 8: Risk Analysis
    create_content_slide(prs, "Risk Analysis", 
                        "• Market risks\n• Operational risks\n• Financial risks\n• Mitigation strategies")
    
    # Slide 9: Implementation Timeline
    create_content_slide(prs, "Implementation Timeline",
                        "• Phase 1: Planning and Setup\n• Phase 2: Implementation\n• Phase 3: Optimization\n• Phase 4: Scale")
    
    # Slide 10: Resource Requirements
    create_content_slide(prs, "Resource Requirements",
                        "• Personnel needs\n• Technology infrastructure\n• Capital requirements\n• Training needs")
    
    # Slide 11: Financial Summary
    if financial_data:
        total_revenue = sum(fd.get('revenue', 0) for fd in financial_data)
        total_net_income = sum(fd.get('net_income', 0) for fd in financial_data)
        summary = f"• Total Revenue (5-year): ${total_revenue:,.0f}\n• Total Net Income: ${total_net_income:,.0f}\n• ROI Analysis"
        create_content_slide(prs, "Financial Summary", summary)
    else:
        create_content_slide(prs, "Financial Summary", "• Financial highlights to be added")
    
    # Slide 12: Conclusion & Recommendations
    create_content_slide(prs, "Conclusion & Recommendations",
                        "• Strategic alignment confirmed\n• Financial viability established\n• Recommended next steps")
    
    # Save to temp file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp_file.name)
    return temp_file.name


@bp.route('/pptx', methods=['POST'])
def export_pptx():
    """Export business case to PowerPoint"""
    try:
        data = request.get_json()
        file_path = generate_business_case_pptx(data)
        project_name = data.get('project_name', 'business_case')
        filename = f"{project_name.replace(' ', '_')}_business_case.pptx"
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        abort(500, description=str(e))


@bp.route('/template', methods=['GET'])
def get_template_info():
    """Get information about the PPTX template structure"""
    return jsonify({
        "slides": [
            {"id": 1, "title": "Title Slide", "type": "title"},
            {"id": 2, "title": "Executive Summary", "type": "content"},
            {"id": 3, "title": "Revenue Projection", "type": "chart"},
            {"id": 4, "title": "Cost Analysis", "type": "chart"},
            {"id": 5, "title": "Profitability Analysis", "type": "chart"},
            {"id": 6, "title": "Net Income Projection", "type": "chart"},
            {"id": 7, "title": "Key Assumptions", "type": "content"},
            {"id": 8, "title": "Risk Analysis", "type": "content"},
            {"id": 9, "title": "Implementation Timeline", "type": "content"},
            {"id": 10, "title": "Resource Requirements", "type": "content"},
            {"id": 11, "title": "Financial Summary", "type": "content"},
            {"id": 12, "title": "Conclusion & Recommendations", "type": "content"}
        ]
    })
